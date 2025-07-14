import sys
from pptx import Presentation
import os

def extract_text_from_pptx(pptx_path, output_path):
    """Extracts all text from a PowerPoint presentation and saves it to a file."""
    try:
        prs = Presentation(pptx_path)
        with open(output_path, 'w', encoding='utf-8') as f:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            f.write(run.text + ' ')
                        f.write('\\n')
                f.write('\\n---\\n\\n') # Slide separator
        print(f"Text extracted successfully to {output_path}")
    except Exception as e:
        print(f"Error extracting text: {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python extract_pptx.py <input_pptx_path> <output_txt_path>", file=sys.stderr)
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    output_file = sys.argv[2]
    
    extract_text_from_pptx(pptx_file, output_file) 